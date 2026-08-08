"""Conversores documento -> Markdown, uno por formato/motor.

Todos comparten la misma firma y devuelven el mismo dict de resultado, para que
la CLI, el lote y el worker de la web los traten igual:

    convert_*(src, out_root, input_root, skip_existing=False, log=None) -> dict

Salida: `<out_root>/<estructura relativa>/<documento>/<documento>.md` con las
imágenes dentro de esa misma carpeta (`assets/`, o `media/` para pandoc, o
planas para Marker) y enlaces relativos al propio `.md`, de modo que la carpeta
de cada documento es autocontenida y portable.
"""
from __future__ import annotations

import os
import re
import shutil
import subprocess
import time
from pathlib import Path

from .paths import doc_out_dir, slugify

# Marker es lento y ocasionalmente se cuelga en un PDF patológico; sin timeout
# un lote nocturno se queda bloqueado en un solo documento. El default es
# generoso a propósito: a ~15-30 s/página, un manual escaneado de varios
# cientos de páginas tarda horas legítimamente y no debe morir por el timeout.
MARKER_TIMEOUT_S = int(os.environ.get("MARKER_TIMEOUT_S", "14400"))
# Vacío = no pasar --mode y dejar que marker elija por dispositivo: `fast`
# (detectores CPU ligeros) en CPU/MPS, `balanced` (layout VLM + OCR de página
# completa) en GPU. Forzar `balanced` en una máquina sin GPU no solo es lento:
# requiere el binario llama-server, y sin él la conversión falla directamente.
# Valores válidos si se fija a mano: "balanced" | "fast".
MARKER_MODE = os.environ.get("MARKER_MODE", "").strip()


def _log(log, msg: str) -> None:
    if log:
        log(msg)


def _result(src, out_dir, md_path, engine, started, **extra) -> dict:
    md_path = Path(md_path) if md_path else None
    text = md_path.read_text(encoding="utf-8", errors="replace") if md_path and md_path.exists() else ""
    return {
        "ok": True,
        "engine": engine,
        "src": str(src),
        "out_dir": str(out_dir),
        "md": str(md_path) if md_path else None,
        "chars": len(text),
        "encabezados": sum(1 for line in text.splitlines() if line.strip().startswith("#")),
        "duracion_s": round(time.time() - started, 1),
        "error": None,
        **extra,
    }


def _error(src, engine, started, error) -> dict:
    return {
        "ok": False, "engine": engine, "src": str(src), "out_dir": None, "md": None,
        "chars": 0, "encabezados": 0, "imagenes": 0,
        "duracion_s": round(time.time() - started, 1), "error": str(error),
    }


def _existing_md(out_dir: Path) -> Path | None:
    """Devuelve el .md ya convertido en out_dir, si lo hay y no está vacío."""
    if not out_dir.is_dir():
        return None
    for md in sorted(out_dir.glob("*.md")):
        if md.stat().st_size > 0:
            return md
    return None


# --------------------------------------------------------------------------- PDF (pymupdf4llm)

def convert_pymupdf(src, out_root, input_root, skip_existing=False, log=None, dpi=150) -> dict:
    """PDF nativo -> Markdown. Rápido, solo CPU, sin modelos."""
    import pymupdf4llm

    started = time.time()
    src = Path(src)
    out_dir = doc_out_dir(src, out_root, input_root)
    try:
        if skip_existing and (md := _existing_md(out_dir)):
            _log(log, f"  SKIP (ya convertido): {src.name}")
            return _result(src, out_dir, md, "pymupdf", started,
                           imagenes=len(list((out_dir / "assets").glob("*"))), skipped=True)

        assets_dir = out_dir / "assets"
        assets_dir.mkdir(parents=True, exist_ok=True)

        md_text = pymupdf4llm.to_markdown(
            str(src),
            write_images=True,
            image_path=str(assets_dir),
            image_format="png",
            dpi=dpi,
        )

        # pymupdf4llm incrusta los enlaces de imagen relativos al CWD, no al .md.
        # Se reescribe cualquier prefijo que termine en assets/ para que el
        # markdown sea portable si se mueve la carpeta o se abre suelto.
        md_text = re.sub(r'(!\[[^\]]*\]\()[^)]*?assets/([^)]+\))', r'\1assets/\2', md_text)

        md_path = out_dir / f"{slugify(src.stem)}.md"
        md_path.write_text(md_text, encoding="utf-8")

        n_images = len(list(assets_dir.glob("*")))
        _log(log, f"  OK ({round(time.time()-started,1)}s): {src.name} | {len(md_text)} chars | {n_images} img")
        return _result(src, out_dir, md_path, "pymupdf", started, imagenes=n_images, skipped=False)
    except Exception as e:
        _log(log, f"  ERROR pymupdf {src.name}: {e}")
        return _error(src, "pymupdf", started, e)


# --------------------------------------------------------------------------- PDF (Marker)

def _flatten_marker_output(out_dir: Path) -> None:
    """Aplana la subcarpeta redundante que crea marker_single.

    marker_single escribe en `<output_dir>/<nombre original>/<archivo>.md` — un
    nivel de más, y con el nombre SIN sanear aunque `--output_dir` sí lo esté.
    """
    for sub in sorted(p for p in out_dir.iterdir() if p.is_dir()):
        if sub.name == "assets" or not any(sub.glob("*.md")):
            continue
        for item in sub.iterdir():
            destino = out_dir / item.name
            if destino.exists():
                shutil.rmtree(destino) if destino.is_dir() else destino.unlink()
            shutil.move(str(item), str(destino))
        sub.rmdir()


def convert_marker(src, out_root, input_root, skip_existing=False, log=None) -> dict:
    """PDF complejo o escaneado -> Markdown vía Marker (layout VLM + OCR completo).

    Mucho más lento que pymupdf4llm pero resuelve multicolumna, jerarquía de
    encabezados en documentos densos, y es el único de los dos que hace OCR real.
    """
    started = time.time()
    src = Path(src)
    out_dir = doc_out_dir(src, out_root, input_root)

    if skip_existing and (md := _existing_md(out_dir)):
        _log(log, f"  SKIP (ya convertido): {src.name}")
        return _result(src, out_dir, md, "marker", started,
                       imagenes=len([p for p in out_dir.iterdir() if p.suffix.lower() != ".md"]),
                       skipped=True)

    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = ["marker_single", str(src), "--output_dir", str(out_dir)]
    if MARKER_MODE:
        cmd += ["--mode", MARKER_MODE]
    _log(log, f"  marker_single (modo: {MARKER_MODE or 'auto por dispositivo'}): "
              f"{src.name} (puede tardar minutos)")
    try:
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=MARKER_TIMEOUT_S)
    except FileNotFoundError:
        return _error(src, "marker", started,
                      "marker_single no está en el PATH (pip install marker-pdf)")
    except subprocess.TimeoutExpired:
        _log(log, f"  ERROR marker {src.name}: timeout tras {MARKER_TIMEOUT_S}s")
        return _error(src, "marker", started, f"timeout tras {MARKER_TIMEOUT_S}s")

    if proc.returncode != 0:
        cola = (proc.stderr or proc.stdout or "").strip().splitlines()[-5:]
        _log(log, f"  ERROR marker {src.name}: rc={proc.returncode} {' | '.join(cola)}")
        return _error(src, "marker", started, f"marker_single rc={proc.returncode}: {' '.join(cola)}")

    try:
        _flatten_marker_output(out_dir)
    except Exception as e:
        _log(log, f"  AVISO: no se pudo aplanar la salida de marker en {out_dir}: {e}")

    md = _existing_md(out_dir)
    if md is None:
        return _error(src, "marker", started, "marker terminó sin producir un .md no vacío")

    n_images = len([p for p in out_dir.iterdir() if p.is_file() and p.suffix.lower() != ".md"])
    _log(log, f"  OK ({round(time.time()-started,1)}s): {src.name} | {n_images} img")
    return _result(src, out_dir, md, "marker", started, imagenes=n_images, skipped=False)


# --------------------------------------------------------------------------- DOCX (pandoc)

def convert_docx(src, out_root, input_root, skip_existing=False, log=None) -> dict:
    """DOCX -> Markdown (GFM) con pandoc, extrayendo las imágenes embebidas.

    pandoc solo extrae imágenes referenciadas en el cuerpo (no encabezados ni
    pies): si `--extract-media` deja imágenes que no aparecen enlazadas en el
    .md, son decorativas.
    """
    started = time.time()
    src = Path(src)
    out_dir = doc_out_dir(src, out_root, input_root)

    if skip_existing and (md := _existing_md(out_dir)):
        _log(log, f"  SKIP (ya convertido): {src.name}")
        return _result(src, out_dir, md, "pandoc", started,
                       imagenes=len(list((out_dir / "media").rglob("*"))), skipped=True)

    out_dir.mkdir(parents=True, exist_ok=True)
    md_name = f"{slugify(src.stem)}.md"
    md_path = out_dir / md_name
    # Se ejecuta pandoc DENTRO de out_dir con --extract-media=. : si se le pasa
    # una ruta absoluta, pandoc escribe esa ruta absoluta dentro del .md
    # (`<img src="/data/.../media/x.png">`) y la carpeta deja de ser portable —
    # muévela y las imágenes se rompen.
    cmd = ["pandoc", str(src.resolve()), "-f", "docx", "-t", "gfm",
           "--extract-media=.", "-o", md_name]
    try:
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=600,
                              cwd=str(out_dir))
    except FileNotFoundError:
        return _error(src, "pandoc", started, "pandoc no está instalado (brew install pandoc)")
    except subprocess.TimeoutExpired:
        return _error(src, "pandoc", started, "timeout de pandoc (600s)")

    if proc.returncode != 0 or not md_path.exists():
        cola = (proc.stderr or "").strip().splitlines()[-3:]
        _log(log, f"  ERROR pandoc {src.name}: {' | '.join(cola)}")
        return _error(src, "pandoc", started, f"pandoc rc={proc.returncode}: {' '.join(cola)}")

    # Cinturón y tirantes: normaliza cualquier ruta que haya quedado absoluta o
    # con el prefijo "./", en las dos sintaxis que emite pandoc.
    text = md_path.read_text(encoding="utf-8", errors="replace")
    abs_prefix = re.escape(str(out_dir.resolve())) + r'/+'
    text = re.sub(abs_prefix, "", text)
    text = re.sub(r'(<img\b[^>]*?\bsrc\s*=\s*["\'])\./', r'\1', text, flags=re.IGNORECASE | re.DOTALL)
    text = re.sub(r'(!\[[^\]]*\]\()\./', r'\1', text)
    md_path.write_text(text, encoding="utf-8")

    n_images = len([p for p in (out_dir / "media").rglob("*") if p.is_file()]) if (out_dir / "media").exists() else 0
    _log(log, f"  OK ({round(time.time()-started,1)}s): {src.name} | {n_images} img")
    return _result(src, out_dir, md_path, "pandoc", started, imagenes=n_images, skipped=False)


# --------------------------------------------------------------------------- PPTX

# Formatos que un visor de Markdown puede renderizar. Los demás (emf/wmf, típicos
# de diagramas pegados desde Office) se extraen igual pero se enlazan como
# adjunto, no como imagen: un ![](x.emf) sale roto en cualquier visor y además
# ensucia el control de calidad con falsos enlaces rotos.
PPTX_RENDERABLE = {"png", "jpg", "jpeg", "gif", "webp", "bmp"}


def _escape_cell(text: str) -> str:
    """Una celda con | o saltos de línea rompe la tabla Markdown entera."""
    return text.replace("|", "\\|").replace("\r", " ").replace("\n", "<br>").strip()


def _iter_shapes(shapes):
    """Recorre las formas recursivamente: los grupos anidan contenido real.

    Sin esto se pierde en silencio todo el texto e imágenes dentro de un grupo,
    que en presentaciones reales es la mayoría del contenido.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def convert_pptx(src, out_root, input_root, skip_existing=False, log=None) -> dict:
    """PPTX -> Markdown: texto por diapositiva, tablas, notas e imágenes extraídas."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    started = time.time()
    src = Path(src)
    out_dir = doc_out_dir(src, out_root, input_root)

    if skip_existing and (md := _existing_md(out_dir)):
        _log(log, f"  SKIP (ya convertido): {src.name}")
        return _result(src, out_dir, md, "pptx", started,
                       imagenes=len(list((out_dir / "assets").glob("*"))), skipped=True)

    try:
        assets_dir = out_dir / "assets"
        assets_dir.mkdir(parents=True, exist_ok=True)

        prs = Presentation(str(src))
        md_lines = [f"# {src.stem}\n"]
        img_count = 0

        for i, slide in enumerate(prs.slides, start=1):
            md_lines.append(f"## Diapositiva {i}\n")
            title_shape = slide.shapes.title

            for shape in _iter_shapes(slide.shapes):
                if title_shape is not None and shape == title_shape:
                    if shape.text.strip():
                        md_lines.append(f"**{shape.text.strip()}**\n")
                    continue

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_count += 1
                    ext = (shape.image.ext or "bin").lower()
                    img_name = f"slide{i:03d}-{img_count:03d}.{ext}"
                    (assets_dir / img_name).write_bytes(shape.image.blob)
                    if ext in PPTX_RENDERABLE:
                        md_lines.append(f"![](assets/{img_name})\n")
                    else:
                        md_lines.append(f"[archivo adjunto: {img_name}](assets/{img_name})\n")

                elif shape.has_table:
                    rows = [[_escape_cell(c.text) for c in row.cells] for row in shape.table.rows]
                    if rows:
                        md_lines.append("| " + " | ".join(rows[0]) + " |")
                        md_lines.append("|" + "---|" * len(rows[0]))
                        for r in rows[1:]:
                            md_lines.append("| " + " | ".join(r) + " |")
                        md_lines.append("")

                elif shape.has_text_frame:
                    # para.text (no join de runs): un párrafo sin runs explícitos
                    # todavía puede tener texto.
                    lines = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                    for line in lines:
                        md_lines.append(f"- {line}")
                    if lines:
                        md_lines.append("")

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    md_lines.append(f"> Notas: {notes}\n")
            md_lines.append("")

        md_text = "\n".join(md_lines)
        md_path = out_dir / f"{slugify(src.stem)}.md"
        md_path.write_text(md_text, encoding="utf-8")

        _log(log, f"  OK ({round(time.time()-started,1)}s): {src.name} | {len(prs.slides)} diapositivas | {img_count} img")
        return _result(src, out_dir, md_path, "pptx", started, imagenes=img_count, skipped=False)
    except Exception as e:
        _log(log, f"  ERROR pptx {src.name}: {e}")
        return _error(src, "pptx", started, e)


# --------------------------------------------------------------------------- despacho

CONVERTERS = {
    "pymupdf": convert_pymupdf,
    "marker": convert_marker,
    "pandoc": convert_docx,
    "pptx": convert_pptx,
}


def engine_for(src, pdf_engine="pymupdf") -> str:
    ext = Path(src).suffix.lower()
    if ext == ".pdf":
        return pdf_engine
    if ext == ".docx":
        return "pandoc"
    if ext == ".pptx":
        return "pptx"
    raise ValueError(f"Formato no soportado: {ext}")


def convert(src, out_root, input_root, pdf_engine="pymupdf", skip_existing=False, log=None) -> dict:
    engine = engine_for(src, pdf_engine)
    return CONVERTERS[engine](src, out_root, input_root, skip_existing=skip_existing, log=log)


def convert_batch(sources, out_root, input_root=None, pdf_engine="pymupdf",
                  skip_existing=False, log=None, on_result=None) -> list[dict]:
    """Convierte un lote. Un fallo NO detiene el resto: se registra y se sigue."""
    from .paths import common_root

    sources = [Path(s) for s in sources]
    input_root = Path(input_root) if input_root else common_root(sources)
    results = []
    for n, src in enumerate(sources, start=1):
        _log(log, f"[{n}/{len(sources)}] {src.name}")
        try:
            r = convert(src, out_root, input_root, pdf_engine=pdf_engine,
                        skip_existing=skip_existing, log=log)
        except ValueError as e:  # formato no soportado
            r = _error(src, "?", time.time(), e)
            _log(log, f"  ERROR: {e}")
        results.append(r)
        if on_result:
            on_result(r)
    return results
